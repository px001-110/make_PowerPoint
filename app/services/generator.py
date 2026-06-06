# -*-coding: UTF-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document


def make_powerpoint(path, prs):

    doc = Document(path)

    # 段落取得
    lines = [p.text.strip() for p in doc.paragraphs]

    if not lines:
        raise ValueError("Wordファイルに内容がありません")
    
    #1行目をタイトルとする
    title_text = lines[0]

    # スライドを生成し、１枚目を追加する。
    title_slide_layout = prs.slide_layouts[0]  # スライドのレイアウトの読み込み
    slide = prs.slides.add_slide(title_slide_layout)  # スライドの追加

    title_placeholder = slide.placeholders[0]  # タイトルのプレースホルダーを取得

    title_placeholder.text = title_text  # タイトルのテキストを設定


    ### 2行目以降を空行区切り
    
    groups = []
    current_group = []

    slide_layout = prs.slide_layouts[1]

    for line in lines[1:]:
        if line == "":
            if current_group:
                groups.append(current_group)
                current_group = []
        else:
            current_group.append(line)
    
    if current_group:
        groups.append(current_group)

    ### グループごとにスライド作成
    for group in groups:
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.placeholders[0]

        body_shape = slide.placeholders[1]


        for i, line in enumerate(group):
            if i == 0:
                title_shape.text = line
            else:
                if body_shape.text_frame.text == "":
                    body_shape.text_frame.text = line
                else:
                    p = body_shape.text_frame.add_paragraph()
                    p.text = line

