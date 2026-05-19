# -*-coding: UTF-8 -*-
from typing_extensions import Self
import wx
from pptx import Presentation
import docx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
import copy


class FileDropTarget(wx.FileDropTarget):
    # Drag & Drop Class
    def __init__(self, window):
        wx.FileDropTarget.__init__(self)
        self.window = window  # ファイルをドロップする対象

    def OnDropFiles(self, x, y, files):  # ファイルをドロップするときの処理
        # self.window.text_entry.SetLabel(files[0])
        for file in files:
            self.window.SomeFunc(file)
        return 0


class App(wx.Frame):
    ### GUI ###
    def __init__(self, parent, id, title):
        wx.Frame.__init__(self, parent, id, title, size=(
            500, 200), style=wx.DEFAULT_FRAME_STYLE)

        # パネル
        p = wx.Panel(self, wx.ID_ANY)

        label = wx.StaticText(p, wx.ID_ANY, 'ここにファイルをドロップしてください',
                              style=wx.SIMPLE_BORDER | wx.TE_CENTER)

        # ドロップ対象の設定
        dt1 = FileDropTarget(self)  # ドロップする対象をフレームに割り当てる
        label.SetDropTarget(dt1)

        # テキスト入力ウィジット
        self.text_entry = wx.TextCtrl(p, wx.ID_ANY)

        # レイアウト
        layout = wx.BoxSizer(wx.VERTICAL)
        layout.Add(label, flag=wx.EXPAND | wx.ALL, border=10, proportion=1)
        layout.Add(self.text_entry, flag=wx.EXPAND | wx.ALL, border=10)
        p.SetSizer(layout)

        self.Show()

    def SomeFunc(self, path):
        self.text_entry.SetLabel(path)
        print(path)
        prs = Presentation(file_path)

        if 'pptx' in path:
            try:
                copy_powerpoint.move_slide(path, prs)
            except:
                pass

        elif 'docx' in path:
            try:
                self.make_pp(path, prs)
            except:
                pass

        prs.save(file_path)
        return

    def make_pp(self, folder, prsen):
        text_list = utility.load_word(folder)
        new_slide_flag = True

        # ここからパワポ
        # インスタンスを生成
        prs = prsen
        # スライドを生成し、１枚目を追加する。
        title_slide_layout = prs.slide_layouts[6]  # スライドのレイアウトの読み込み
        slide = prs.slides.add_slide(title_slide_layout)  # スライドの追加
        txbox_left = utility.Centis(0.5)
        txbox_top = utility.Centis(0.8)
        # タイトル用のテキストボックスの作成

        # テキストボックスの幅と高さを決める
        txbox_width = utility.Centis(24)
        txbox_height = utility.Centis(17)
        txBox = slide.shapes.add_textbox(
            txbox_left, txbox_top, txbox_width, txbox_height)

        for txt in text_list:
            if txt == "":
                title_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(title_slide_layout)
                txbox_left = utility.Centis(0.5)
                txbox_top = utility.Centis(0.8)

                # テキストボックスの幅と高さを決める
                txbox_width = utility.Centis(24)
                txbox_height = utility.Centis(2)
                txBox = slide.shapes.add_textbox(
                    txbox_left, txbox_top, txbox_width, txbox_height)
                new_slide_flag = True

            elif new_slide_flag:
                tf = txBox.text_frame
                tf.word_wrap = True
                txBox.text = txt
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(44)
                    paragraph.font.bold = True
                    paragraph.font.name = 'HGP創英ﾌﾟﾚｾﾞﾝｽEB'

                new_slide_flag = False

            else:
                tf = txBox.text_frame
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(40)
                    paragraph.font.bold = True
                    paragraph.font.name = 'MS PGothic'
                    paragraph.line_spacing = 1.5
                tf.word_wrap = True
                p = tf.add_paragraph()
                utility.add_text(p, txt, 40, True, 'MS PGothic')


class utility():
    def Centis(length):
        centi = Inches(length/2.54)
        return centi
# テキストボックスを作り、文字を入れる

    def make_textbox(font_size, txt, slide):
        # テキストボックスの位置を決める
        textbox_left = Self.Centis(0.5)
        textbox_top = Self.Centis(0.8)
        # テキストボックスの幅と高さを決める
        textbox_width = Self.Centis(24)
        textbox_height = Self.Centis(17)
        # テキストのフォントサイズを決める　
        # テキストボックスを配置する
        textbox = slide.shapes.add_textbox(
            textbox_left, textbox_top, textbox_width, textbox_height)
        # テキストボックスに書き込む
        textbox.text = txt
        for paragraph in textbox.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = 'MS PGothic'
            paragraph.font.bold = True
            paragraph.line_spacing = 1.5
        # テキストの自動折り返しができそう
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        # text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def add_text(p, msg, font_size, is_bold, f_name):
        p.text = msg
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.name = f_name
        p.line_spacing = 1.5

    def load_word(folder):
        file = folder
        text_list = []
        doc = docx.Document(file)
        for paragraph in doc.paragraphs:
            text_list.append(paragraph.text)

        return text_list


class copy_powerpoint():
    def _get_blank_slide_layout(pres):
        layout_items_count = [len(layout.placeholders)
                              for layout in pres.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return pres.slide_layouts[blank_layout_id]

    def move_slide(path, pres):
        pres1 = Presentation(path)
        for slide in pres1.slides:
            # source = pres1.slides[index]
            source = slide
            blank_slide_layout = copy_powerpoint._get_blank_slide_layout(pres)
            dest = pres.slides.add_slide(blank_slide_layout)

            for shape in source.shapes:
                newel = copy.deepcopy(shape.element)
                dest.shapes._spTree.insert_element_before(newel, 'p:extLst')


app = wx.App()
dialog = wx.FileDialog(None, u'ファイル名を入力してください。', wildcard="pptx files(*.pptx)|*.pptx",
                       style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)
if dialog.ShowModal() == wx.ID_OK:
    file_path = dialog.GetPath()
    prs = Presentation()
    prs.save(file_path)
    App(None, -1, '自動生成')
    app.MainLoop()
elif dialog.ShowModal == wx.ID_CANCEL:
    dialog.Destroy()
    exit()
