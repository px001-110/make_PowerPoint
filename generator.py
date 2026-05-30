# -*-coding: UTF-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
# import re
from docx import Document
from copy import deepcopy
from pptx.enum.text import PP_ALIGN

class PowerPointGenerator:
           
    def make_PowerPoint(self, path, prs):

        doc = Document(path)

        # 段落取得
        lines = [p.text.strip() for p in doc.paragraphs]

        if not lines:
            raise ValueError("Wordファイルに内容がありません")
        
        #1行目をタイトルとする
        title_text = lines[0]

        # スライドを生成し、１枚目を追加する。
        title_slide_layout = prs.slide_layouts[0]  # スライドのレイアウトの読み込み
        slide = prs.slides.add_slide(title_slide_layout)  # スライドの追加

        title_placeholder = slide.placeholders[0]  # タイトルのプレースホルダーを取得

        title_placeholder.text = title_text  # タイトルのテキストを設定


        ### 2行目以降を空行区切り
        
        groups = []
        current_group = []

        slide_layout = prs.slide_layouts[1]

        for line in lines[1:]:
            if line == "":
                if current_group:
                    groups.append(current_group)
                    current_group = []
            else:
                current_group.append(line)
        
        if current_group:
            groups.append(current_group)
        
        print(groups)  # グループ分けの確認

        ### グループごとにスライド作成
        for group in groups:
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.placeholders[0]

            body_shape = slide.placeholders[1]

            # txBox = slide.shapes.add_textbox(
            #     Utility.centis(0.5),
            #     Utility.centis(0.5),
            #     Utility.centis(30),
            #     Utility.centis(20)
            # )
            # text_frame = txBox.text_frame
            # text_frame.clear()

            for i, line in enumerate(group):
                if i == 0:
                    title_shape.text = line
                else:
                    if body_shape.text_frame.text == "":
                        body_shape.text_frame.text = line
                    else:
                        p = body_shape.text_frame.add_paragraph()
                        p.text = line

                # p.text = line
                # p.font.size = Pt(54)
                # p.font.bold = True
                # p.font.name = "MS PGothic"
                # p.line_spacing = Pt(84)
                # p.alignment = PP_ALIGN.CENTER

            # txBox.top = Inches(2)
    



# 修正しようとしたやつ
    def move_slide(self, path, pres):
        
        src = Presentation(path)

        for slide in src.slides:

            blank_layout = pres.slide_layouts[6]

            new_slide = pres.slides.add_slide(blank_layout)

            for shape in slide.shapes:

                el = deepcopy(shape.element)

                new_slide.shapes._spTree.insert_element_before(
                    el,
                    'p:extLst'
                )
        
        # alignment 修正
            for old_shape, new_shape in zip(slide.shapes, new_slide.shapes):

                if not old_shape.has_text_frame:
                    continue

                for old_p, new_p in zip(
                    old_shape.text_frame.paragraphs,
                    new_shape.text_frame.paragraphs
                ):
                    
                    new_p.alignment = old_p.alignment

            # # 背景コピー
            # if slide.background:
            #     new_slide.background.fill.fore_color.rgb = \
            #         slide.background.fill.for_color.rgb



class Utility:

    @staticmethod
    def centis(length):
        return Inches(length/2.54)

# テキストボックスを作り、文字を入れる
    @staticmethod
    def add_text(p, msg, font_size, is_bold, font_name):
        p.text = msg
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.name = font_name
        p.line_spacing = 1.5

# class WordReader:

#     @staticmethod
#     def read_docx(path):

#         doc = Document(path)
        
#         texts = []
#         for para in doc.paragraphs:

#             text = para.text.strip()
            
#             if text:
#                 texts.append(text)

#         return texts


class copy_powerpoint():
    def _get_blank_slide_layout(pres):
        layout_items_count = [len(layout.placeholders)
                              for layout in pres.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return pres.slide_layouts[blank_layout_id]